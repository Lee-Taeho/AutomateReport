import pptx.text.text
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
# from scrapper import *
import pandas as pd

year = 2022
month = 11
# bot_name = "Astro"
bot_name = "Nate & Maven"
month_to_name = {
    1: "Jan", 2: "Feb", 3: "Mar",
    4: "Apr", 5: "May", 6: "Jun",
    7: "Jul", 8: "Aug", 9: "Sep",
    10: "Oct", 11: "Nov", 12: "Dec"
}

# page 3
user_and_messages_img_path = "web-element-2.png"
messages_img_path = "web-element-3.png"

# page 4
conversations_img_path = "web-element-4.png"
irqa_img_path = "web-element-5.png"




def update_date(text: pptx.text.text.TextFrame, month: str, year):
    text.word_wrap = False
    p = text.paragraphs[0]
    p.clear()

    run = text.paragraphs[0].add_run()
    run.text = bot_name
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)
    font.color.rgb = RGBColor(8, 172, 84)

    run = text.paragraphs[0].add_run()
    run.text = f" {month} {year}"
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(16)


ppt = Presentation("resource/report.pptx")


# 2nd page
# slide = ppt.slides[1]
# filename = get_excel_file_name()
# df = pd.read_excel(filename, sheet_name="Users and mesages")
# print(df)

# 3rd page
slide = ppt.slides[2]
texts = []
for shape in slide.shapes:
    if shape.shape_type == 17:
        texts.append(shape)
    elif shape.shape_type == 13:
        slide.shapes.element.remove(shape.element)

# Freeform 5: page number
# placeholder 14: Talkbot Report
# Text_Box 17: Nate & Maven Feb 2022
# Picture 13: picture

# update date for the 3rd page
date = texts[0].text_frame
update_date(date, month_to_name[month], year)

# update graph for the 3rd page
users_and_messages = slide.shapes.add_picture(user_and_messages_img_path, Inches(0.8), Inches(1.2), width=Inches(10.1))
messages = slide.shapes.add_picture(messages_img_path, Inches(0.8), Inches(1.2) + users_and_messages.height,
                                    width=Inches(7))


# 4th page
slide = ppt.slides[3]


texts = []
for shape in slide.shapes:
    if shape.shape_type == 17:
        texts.append(shape)
    elif shape.shape_type == 13:
        slide.shapes.element.remove(shape.element)
    print(shape.shape_type)

conversations = slide.shapes.add_picture(conversations_img_path, Inches(0.8), Inches(1.2), width=Inches(7))
irqa = slide.shapes.add_picture(irqa_img_path, Inches(0.8), Inches(1.2) + conversations.height, width=Inches(7))

text = texts[0].text_frame
update_date(text, month_to_name[month], year)


# 5th page
slide = ppt.slides[4]

print("\n")
texts = []
for shape in slide.shapes:
    if shape.shape_type == 17:
        texts.append(shape)
    elif shape.shape_type == 13:
        slide.shapes.element.remove(shape.element)
    print(shape.shape_type)

# text_box 0 : page number
# text_box 1: Nate & Maven Feb 2022



text = texts[1].text_frame

emails = slide.shapes.add_picture("email.png", Inches(0.8), Inches(1.2), width=Inches(10))
update_date(text, month_to_name[month], year)



ppt.save('New_report.pptx')
